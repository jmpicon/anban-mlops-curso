"""
PPTX Día 2 — MLOps en producción (ANBAN).
Reusa los helpers de build_dia1 (mismo estilo).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
CYAN = RGBColor(0x00, 0xB4, 0xD8)
CORAL = RGBColor(0xFF, 0x6B, 0x35)
BG = RGBColor(0xF8, 0xFA, 0xFC)
INK = RGBColor(0x1A, 0x1D, 0x29)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0xA0, 0x55, 0xF7)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PINK = RGBColor(0xEC, 0x49, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False
    return bg

def add_text(slide, text, left, top, width, height, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color
    return tb

def add_rect(slide, l, t, w, h, *, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s

def add_round(slide, l, t, w, h, *, fill=CYAN, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s

def header(slide, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), fill=NAVY)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.06), fill=CYAN)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(11), Inches(0.6),
             size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        add_text(slide, kicker, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.45),
                 size=12, color=CYAN, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, bold=True)
    add_text(slide, "MLOps / DataOps · ANBAN · José Picón",
             Inches(0.5), Inches(7.05), Inches(8), Inches(0.3), size=10, color=GREY)

def footer_pagenum(slide, n, total):
    add_text(slide, f"{n} / {total}",
             Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             size=10, color=GREY, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, h, *, size=18, color=INK, bullet="●"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(8)
        if isinstance(item, tuple):
            head, sub = item
            r = p.add_run(); r.text = f"{bullet}  "
            r.font.size = Pt(size); r.font.color.rgb = CYAN; r.font.bold = True
            r2 = p.add_run(); r2.text = head
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.bold = True
            r3 = p.add_run(); r3.text = f" — {sub}"
            r3.font.size = Pt(size-2); r3.font.color.rgb = GREY
        else:
            r = p.add_run(); r.text = f"{bullet}  "
            r.font.size = Pt(size); r.font.color.rgb = CYAN; r.font.bold = True
            r2 = p.add_run(); r2.text = item
            r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

# === PORTADA & MISC ===
def slide_cover():
    s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.6), SH, fill=CORAL)
    add_text(s, "MLOps / DataOps", Inches(1.0), Inches(2.0), Inches(11), Inches(1.4),
             size=60, bold=True, color=WHITE)
    add_text(s, "Producción: serving, CI/CD y monitoring",
             Inches(1.0), Inches(3.2), Inches(11), Inches(0.8),
             size=28, color=CORAL)
    add_rect(s, Inches(1.0), Inches(4.1), Inches(3), Inches(0.05), fill=CYAN)
    add_text(s, "Día 2 — De Staging a Production y vivir contándolo",
             Inches(1.0), Inches(4.3), Inches(11), Inches(0.6),
             size=22, color=WHITE)
    add_text(s, "ANBAN · Asociación Big Data España",
             Inches(1.0), Inches(6.0), Inches(11), Inches(0.5),
             size=16, color=CORAL)
    add_text(s, "José Picón · 2026",
             Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             size=14, color=WHITE)

def slide_d1_recap():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Recap Día 1", "intro")
    bullets(s, [
        "Diagnosticamos madurez MLOps (L0/L1/L2) y mapamos roles del equipo.",
        "DataOps = pipelines + versionado (DVC) + calidad (GE) + lineage.",
        "MLflow nos da Tracking + Models + Registry para experimentar y promocionar.",
        "Tenemos un modelo en Staging listo para ser servido.",
    ], Inches(0.5), Inches(1.4), Inches(12.3), Inches(3), size=20)
    add_round(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), fill=CORAL)
    add_text(s, "Hoy: cómo lo llevamos a producción y cómo lo mantenemos vivo",
             Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_agenda_d2():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Agenda · Día 2 (4 horas)", "agenda")
    rows = [
        ("00:00–01:00", "M4", "Empaquetado y serving (FastAPI + Docker)", CYAN),
        ("01:00–01:15", "—",  "Pausa", GREY),
        ("01:15–02:30", "M5", "CI/CD para ML + Monitoring + Drift", CORAL),
        ("02:30–02:45", "—",  "Pausa", GREY),
        ("02:45–03:45", "M6", "Caso end-to-end + Governance", GREEN),
        ("03:45–04:00", "—",  "Cierre y siguientes pasos", NAVY),
    ]
    y = Inches(1.4)
    for h, code, txt, col in rows:
        add_round(s, Inches(0.5), y, Inches(1.7), Inches(0.7), fill=col)
        add_text(s, h, Inches(0.5), y, Inches(1.7), Inches(0.7),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(2.4), y, Inches(0.9), Inches(0.7), fill=NAVY)
        add_text(s, code, Inches(2.4), y, Inches(0.9), Inches(0.7),
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, txt, Inches(3.5), y, Inches(9.5), Inches(0.7),
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.85)

def slide_module_divider(num, title, subtitle, color):
    s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
    add_rect(s, 0, Inches(2.5), SW, Inches(2.5), fill=color)
    add_text(s, f"Módulo {num}", Inches(0.8), Inches(2.7), Inches(12), Inches(0.6),
             size=22, color=WHITE, bold=True)
    add_text(s, title, Inches(0.8), Inches(3.2), Inches(12), Inches(1.1),
             size=44, color=WHITE, bold=True)
    add_text(s, subtitle, Inches(0.8), Inches(5.2), Inches(12), Inches(0.6),
             size=20, color=CYAN)

# === MÓDULO 4 ===
def slide_serving_patterns():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Patrones de serving", "M4")
    rows = [
        ("Batch",     "Predicciones programadas, escritas a tabla", "Recomendaciones diarias, scoring nocturno", CYAN),
        ("Online",    "API síncrona <100ms (REST/gRPC)",            "Fraude, recommendation realtime, antispam", CORAL),
        ("Streaming", "Predicciones por evento (Kafka/Pub-Sub)",     "IoT, click stream, anomalías",              GREEN),
        ("Edge",      "Modelo en el dispositivo (móvil, browser, IoT)","Latencia mínima, datos sensibles, offline", NAVY),
    ]
    for i, (n, d, u, c) in enumerate(rows):
        y = Inches(1.3 + i*1.4)
        add_round(s, Inches(0.5), y, Inches(2), Inches(1.2), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2), Inches(1.2),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(2.6), y, Inches(10.2), Inches(1.2), fill=WHITE, line=c)
        add_text(s, d, Inches(2.8), y+Inches(0.1), Inches(10), Inches(0.5),
                 size=14, color=INK)
        add_text(s, f"Cuándo: {u}", Inches(2.8), y+Inches(0.65), Inches(10), Inches(0.5),
                 size=13, color=c, bold=True)

def slide_packaging_options():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Empaquetado: cómo se va el modelo de viaje", "M4")
    rows = [
        ("Pickle/Joblib", "Sencillo pero atado a Python y a versión exacta de libs", RED),
        ("ONNX",          "Estándar abierto, multi-framework, multi-runtime",        GREEN),
        ("MLflow flavor", "Modelo + entorno + signature en un único bundle",         CYAN),
        ("BentoML",       "Empaqueta servicio + modelo + Dockerfile autogenerado",   CORAL),
        ("TF Serving",    "Servidor optimizado para SavedModel TF",                  PURPLE),
        ("TorchServe",    "Equivalente para PyTorch, batching dinámico",             ORANGE),
        ("Triton",        "NVIDIA, GPU/CPU, multi-modelo, gRPC",                     PINK),
    ]
    for i, (n, d, c) in enumerate(rows):
        y = Inches(1.3 + i*0.75)
        add_round(s, Inches(0.5), y, Inches(2.8), Inches(0.65), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2.8), Inches(0.65),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(3.4), y, Inches(9.4), Inches(0.65), fill=WHITE, line=c)
        add_text(s, d, Inches(3.6), y, Inches(9.2), Inches(0.65),
                 size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)

def slide_api_design():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Diseñar la API: contrato y validación", "M4")
    add_text(s, "Endpoints mínimos",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        ("POST /predict", "predicción (JSON in, JSON out)"),
        ("GET /health",   "liveness + readiness"),
        ("GET /version",  "modelo, run_id, signature"),
        ("GET /metrics",  "Prometheus expose"),
        ("POST /explain", "shap values (opcional)"),
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(3), size=14)

    add_text(s, "Validación con Pydantic",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    add_round(s, Inches(7), Inches(1.7), Inches(6), Inches(5.2), fill=NAVY)
    code = ("from pydantic import BaseModel, Field\n\n"
            "class Features(BaseModel):\n"
            "    age: int = Field(ge=0, le=120)\n"
            "    workclass: str\n"
            "    education_num: int = Field(ge=0, le=20)\n"
            "    hours_per_week: int = Field(ge=0, le=99)\n"
            "    capital_gain: float = Field(ge=0)\n\n"
            "class Prediction(BaseModel):\n"
            "    label: int\n"
            "    proba: float\n"
            "    model_version: str\n"
            "    request_id: str")
    add_text(s, code, Inches(7.2), Inches(1.85), Inches(5.7), Inches(4.9),
             size=13, color=CYAN, font="Consolas")

def slide_perf_metrics():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Performance: qué medir y cómo", "M4")
    add_text(s, "Métricas de servicio",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        ("Latencia p50/p95/p99", "el promedio engaña, el p99 es el SLA real"),
        ("Throughput RPS", "requests por segundo sostenidos"),
        ("Tasa de error", "5xx + timeouts + payload inválido"),
        ("Saturación", "CPU, memoria, GPU util, queue depth"),
        ("Cold start", "tiempo desde ↑ contenedor hasta primer 200"),
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(4), size=14)

    add_text(s, "Estrategias para bajar latencia",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Quantización (int8/fp16) o destilación.",
        "Batching dinámico en serving.",
        "Cache de predicciones para inputs frecuentes.",
        "Modelos en ONNX runtime (>2× más rápido a veces).",
        "Pre-cargar features desde el feature store.",
        "Async I/O en FastAPI + uvicorn workers.",
    ], Inches(7.2), Inches(1.7), Inches(6), Inches(4), size=14)

def slide_train_serve_skew():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Train ≠ Serve: el bug clásico", "M4")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0), fill=CORAL)
    add_text(s, "El modelo predice peor en producción que en evaluación. ¿Por qué?",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    causes = [
        ("Encoding distinto", "OneHot fitted en train, otro fitted en serve"),
        ("Imputación incoherente", "median train ≠ median real-time"),
        ("Tipos diferentes",       "int en train, str en serve"),
        ("Features faltantes",     "campo opcional en producción que en train siempre venía"),
        ("Tiempo y zona horaria",  "UTC vs local lleva a otro día/semana"),
        ("Versión de librería",    "scikit 1.4 vs 1.5 cambia un default"),
    ]
    for i, (h, t) in enumerate(causes):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(2.5 + row*1.4)
        add_round(s, x, y, Inches(6.2), Inches(1.2), fill=WHITE, line=CYAN)
        add_rect(s, x, y, Inches(0.15), Inches(1.2), fill=CORAL)
        add_text(s, h, x+Inches(0.3), y+Inches(0.1), Inches(5.8), Inches(0.4),
                 size=15, bold=True, color=NAVY)
        add_text(s, t, x+Inches(0.3), y+Inches(0.55), Inches(5.8), Inches(0.6),
                 size=13, color=GREY)
    add_text(s, "Prevención: pipeline serializado completo (sklearn Pipeline / mlflow.pyfunc).",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
             size=15, bold=True, color=NAVY)

def slide_lab3_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Lab 3 — FastAPI + Docker (40 min)", "M4")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7), fill=CYAN)
    add_text(s, "Objetivo: servir el modelo del Lab 2 como API contenedorizada",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Pasos",
             Inches(0.5), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Cargar modelo desde 'models:/income-clf/Staging'.",
        "Implementar /predict, /health, /version con Pydantic.",
        "Dockerfile multi-stage (slim, ~250MB).",
        "Levantar con docker run y curl.",
        "Stress test con Locust 50 usuarios concurrentes.",
        "Validar p95 < 100ms.",
    ], Inches(0.7), Inches(2.6), Inches(6), Inches(4), size=14)
    add_text(s, "Comprobaciones",
             Inches(7), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "GET /health responde 200 < 50ms.",
        "GET /version devuelve run_id y signature.",
        "POST /predict valida input y devuelve label + proba.",
        "Locust report HTML con estadísticas.",
        "Imagen final < 300MB.",
    ], Inches(7.2), Inches(2.6), Inches(6), Inches(4), size=14)

def slide_dockerfile():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Dockerfile multi-stage para serving", "M4")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8), fill=NAVY)
    code = ("# stage 1: builder\n"
            "FROM python:3.11-slim AS builder\n"
            "WORKDIR /build\n"
            "COPY requirements.txt .\n"
            "RUN pip install --user --no-cache-dir -r requirements.txt\n\n"
            "# stage 2: runtime\n"
            "FROM python:3.11-slim\n"
            "WORKDIR /app\n"
            "COPY --from=builder /root/.local /root/.local\n"
            "ENV PATH=/root/.local/bin:$PATH \\\n"
            "    MLFLOW_TRACKING_URI=http://mlflow:5000 \\\n"
            "    MODEL_URI=models:/income-clf/Staging\n"
            "COPY app/ ./app/\n\n"
            "# usuario sin privilegios\n"
            "RUN useradd -u 1000 -m svc && chown -R svc /app\n"
            "USER svc\n\n"
            "EXPOSE 8000\n"
            "HEALTHCHECK --interval=30s --timeout=3s \\\n"
            "  CMD curl -f http://localhost:8000/health || exit 1\n\n"
            "CMD [\"uvicorn\", \"app.main:app\", \"--host\", \"0.0.0.0\", \"--port\", \"8000\"]")
    add_text(s, code, Inches(0.7), Inches(1.4), Inches(12), Inches(5.6),
             size=12, color=CYAN, font="Consolas")

def slide_m4_recap():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Recap Módulo 4", "M4")
    bullets(s, [
        "Elige patrón de serving por requisito de latencia / volumen / contexto.",
        "Empaqueta modelo con su entorno (mlflow flavor) — no solo el .pkl.",
        "API con contrato: Pydantic valida, Swagger documenta gratis.",
        "Mide latencia en p95/p99, no en promedios.",
        "Train/Serve skew: serializa el pipeline completo, no solo el estimador.",
        "Imagen Docker mínima (< 300MB) y user no-root.",
    ], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5), size=18)

# === MÓDULO 5 — CI/CD + Monitoring ===
def slide_ci_cd_ct():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "CI / CD / CT: el triple bucle ML", "M5")
    cols = [
        ("CI", "Continuous Integration",
         "Tests código, tests de datos, lint, build", CYAN),
        ("CD", "Continuous Delivery / Deployment",
         "Build modelo, build imagen, despliegue gradual", CORAL),
        ("CT", "Continuous Training",
         "Reentrenamiento automático por trigger (drift, calendario)", GREEN),
    ]
    for i, (k, n, d, c) in enumerate(cols):
        x = Inches(0.5 + i*4.27)
        add_round(s, x, Inches(1.4), Inches(4.1), Inches(5), fill=c)
        add_text(s, k, x, Inches(1.5), Inches(4.1), Inches(0.8),
                 size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, n, x, Inches(2.5), Inches(4.1), Inches(0.6),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, d, x+Inches(0.3), Inches(3.3), Inches(3.5), Inches(2.5),
                 size=15, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_pipeline_jobs():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Pipeline ML típico: jobs encadenados", "M5")
    jobs = [("Lint+Test", CYAN), ("Validate Data", CORAL), ("Train", GREEN),
            ("Evaluate", NAVY), ("Promote", PURPLE), ("Build Img", ORANGE),
            ("Deploy", PINK)]
    x = Inches(0.4)
    for i, (n, c) in enumerate(jobs):
        add_round(s, x, Inches(2.5), Inches(1.65), Inches(1.0), fill=c)
        add_text(s, n, x, Inches(2.5), Inches(1.65), Inches(1.0),
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(jobs)-1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     x+Inches(1.7), Inches(2.85), Inches(0.12), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background()
        x += Inches(1.82)
    add_text(s, "Detalles por job",
             Inches(0.5), Inches(4.0), Inches(12), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        ("Lint+Test", "ruff + pytest, cobertura, mypy"),
        ("Validate Data", "Great Expectations / Pandera sobre snapshot"),
        ("Train", "dvc repro o python train.py + MLflow log"),
        ("Evaluate", "tests de métrica mínima, fairness, slices"),
        ("Promote", "si supera al modelo Production en ≥ X%, mover stage"),
        ("Build / Deploy", "docker build + push, helm/kubectl rollout"),
    ], Inches(0.7), Inches(4.6), Inches(12), Inches(2.5), size=13)

def slide_github_actions():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "GitHub Actions: workflow ML", "M5")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7), fill=NAVY)
    yaml = ("# .github/workflows/ml.yml\n"
            "name: ml-pipeline\n"
            "on:\n"
            "  pull_request: { branches: [main] }\n"
            "  schedule: [{ cron: '0 4 * * 1' }]\n\n"
            "jobs:\n"
            "  test:\n"
            "    runs-on: ubuntu-latest\n"
            "    steps:\n"
            "      - uses: actions/checkout@v4\n"
            "      - uses: actions/setup-python@v5\n"
            "        with: { python-version: '3.11' }\n"
            "      - run: pip install -r requirements.txt\n"
            "      - run: ruff check . && pytest -v\n\n"
            "  train:\n"
            "    needs: test\n"
            "    runs-on: ubuntu-latest\n"
            "    steps:\n"
            "      - uses: actions/checkout@v4\n"
            "      - run: pip install -r requirements.txt\n"
            "      - run: dvc pull && python src/train.py\n"
            "      - run: python src/promote_if_better.py")
    add_text(s, yaml, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5),
             size=12, color=CYAN, font="Consolas")

def slide_deploy_strategies():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Estrategias de despliegue de modelo", "M5")
    rows = [
        ("Recreate",    "Apaga v1, sube v2. Hay downtime",
         "Cargas internas, no críticas",                 GREY),
        ("Blue/Green",  "Tráfico salta de v1 a v2 entero",
         "Apps con ventana de mantenimiento",            CYAN),
        ("Canary",      "5% → 25% → 50% → 100% con auto-rollback",
         "Producción con tráfico real, riesgo controlado", CORAL),
        ("Shadow",      "v2 recibe copia del tráfico, no responde",
         "Comparar predicciones antes de exponer",        GREEN),
        ("A/B testing", "% de usuarios a v2, comparas KPI negocio",
         "Decisión con métrica de producto",              NAVY),
        ("Multi-armed", "Asignación adaptativa al mejor modelo",
         "Recomendadores, growth experimentation",        PURPLE),
    ]
    for i, (n, d, u, c) in enumerate(rows):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(1.3 + row*1.85)
        add_round(s, x, y, Inches(6.2), Inches(1.6), fill=WHITE, line=c)
        add_rect(s, x, y, Inches(2), Inches(1.6), fill=c)
        add_text(s, n, x, y, Inches(2), Inches(1.6),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, d, x+Inches(2.1), y+Inches(0.15), Inches(4.0), Inches(0.6),
                 size=13, color=INK)
        add_text(s, f"Cuándo: {u}", x+Inches(2.1), y+Inches(0.85), Inches(4.0), Inches(0.6),
                 size=12, color=c, bold=True)

def slide_monitoring_what():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "¿Qué monitorizar en un modelo?", "M5")
    layers = [
        ("Servicio",  "Latencia, throughput, errores, CPU/mem", CYAN),
        ("Datos",     "Distribución de inputs vs referencia, missing rate", CORAL),
        ("Modelo",    "Accuracy proxy, distribución de outputs", GREEN),
        ("Negocio",   "Conversión, ingreso, NPS — el por qué último", NAVY),
    ]
    for i, (n, d, c) in enumerate(layers):
        y = Inches(1.5 + i*1.3)
        add_round(s, Inches(0.5), y, Inches(2.5), Inches(1.0), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2.5), Inches(1.0),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(3.1), y, Inches(9.7), Inches(1.0), fill=WHITE, line=c)
        add_text(s, d, Inches(3.3), y, Inches(9.4), Inches(1.0),
                 size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), fill=CORAL)
    add_text(s, "Sin la capa de negocio, no sabes si el modelo realmente sirve",
             Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_drift_types():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Tipos de drift: P(X), P(y), P(y|X)", "M5")
    rows = [
        ("Data drift\n(covariate shift)", "Cambia P(X). Los inputs entran distinto",
         "Nuevos clientes, estacionalidad, pandemia", CYAN),
        ("Label drift",     "Cambia P(y). La proporción de clases varía",
         "Más fraudes en navidad, menos churn en verano", CORAL),
        ("Concept drift",   "Cambia P(y|X). La relación entrada→salida cambia",
         "Comportamiento de cliente cambia con la inflación", GREEN),
        ("Upstream drift",  "Un sistema upstream cambia formato",
         "Proveedor empieza a mandar 'NA' en lugar de null", NAVY),
    ]
    for i, (n, d, u, c) in enumerate(rows):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(1.3 + row*2.7)
        add_round(s, x, y, Inches(6.2), Inches(2.4), fill=WHITE, line=c)
        add_rect(s, x, y, Inches(6.2), Inches(0.7), fill=c)
        add_text(s, n, x, y, Inches(6.2), Inches(0.7),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, d, x+Inches(0.2), y+Inches(0.85), Inches(5.8), Inches(0.7),
                 size=14, color=INK)
        add_text(s, f"Ejemplo: {u}", x+Inches(0.2), y+Inches(1.55), Inches(5.8), Inches(0.7),
                 size=13, color=c, bold=True)

def slide_drift_tests():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Tests estadísticos para drift", "M5")
    rows = [
        ("PSI (Population Stability Index)", "Compara distribuciones discretas. Umbral típico: <0.1 ok, >0.25 alerta", CYAN),
        ("Kolmogorov-Smirnov",                "Continuo. Sensible al tamaño muestral", CORAL),
        ("Chi-cuadrado",                      "Categórico, frecuencias esperadas vs observadas", GREEN),
        ("Wasserstein / Earth Mover",         "Distancia entre distribuciones, robusto", NAVY),
        ("Jensen-Shannon",                    "Divergencia entre distribuciones, simétrico", PURPLE),
    ]
    for i, (n, d, c) in enumerate(rows):
        y = Inches(1.4 + i*1.0)
        add_round(s, Inches(0.5), y, Inches(3.8), Inches(0.85), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(3.8), Inches(0.85),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(4.4), y, Inches(8.4), Inches(0.85), fill=WHITE, line=c)
        add_text(s, d, Inches(4.6), y, Inches(8.2), Inches(0.85),
                 size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Usa varios tests: cada uno mira algo diferente.",
             Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             size=14, color=GREY, bold=True)

def slide_evidently():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Evidently: drift y model performance", "M5")
    add_text(s, "Qué te da",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Reportes HTML interactivos comparando referencia vs producción.",
        "Tests automatizables (ok/fail) para CI.",
        "Métricas: data drift, target drift, performance.",
        "Servicio de monitoring continuo (Evidently UI).",
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(3), size=14)

    add_text(s, "Código mínimo",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    add_round(s, Inches(7), Inches(1.7), Inches(6), Inches(5.2), fill=NAVY)
    code = ("from evidently.report import Report\n"
            "from evidently.metrics import (\n"
            "    DataDriftPreset,\n"
            "    DataQualityPreset,\n"
            ")\n\n"
            "report = Report(metrics=[\n"
            "    DataDriftPreset(),\n"
            "    DataQualityPreset(),\n"
            "])\n"
            "report.run(\n"
            "    reference_data=ref_df,\n"
            "    current_data=prod_df,\n"
            ")\n"
            "report.save_html('drift.html')\n\n"
            "tests = TestSuite(tests=[NoTargetPerformanceTestPreset()])\n"
            "tests.run(reference_data=ref, current_data=prod)\n"
            "assert tests.as_dict()['summary']['all_passed']")
    add_text(s, code, Inches(7.2), Inches(1.85), Inches(5.7), Inches(4.9),
             size=12, color=CYAN, font="Consolas")

def slide_alerting():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "De la métrica a la acción", "M5")
    add_text(s, "Cuándo dispara una alerta no es trivial",
             Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             size=18, bold=True, color=CORAL)
    bullets(s, [
        ("Threshold simple", "drift > 0.25 → alerta. Fácil pero ruidoso"),
        ("Ventana móvil",   "promedio últimas 7 ejecuciones > umbral"),
        ("CUSUM",           "detecta cambios sostenidos vs picos"),
        ("Anomaly score",   "modelo de anomalías sobre las propias métricas"),
        ("Mixto",           "métrica del modelo + métrica de negocio + tiempo"),
    ], Inches(0.7), Inches(1.9), Inches(12), Inches(3), size=15)

    add_text(s, "Posibles reacciones",
             Inches(0.5), Inches(4.7), Inches(12), Inches(0.5),
             size=18, bold=True, color=GREEN)
    bullets(s, [
        "Avisar al owner (Slack, PagerDuty).",
        "Disparar reentrenamiento (CT job).",
        "Hacer rollback automático al modelo anterior.",
        "Enrutar tráfico a fallback (heurística simple).",
    ], Inches(0.7), Inches(5.3), Inches(12), Inches(2), size=15)

def slide_lab4_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Lab 4 — CI/CD + Monitoring (50 min)", "M5")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7), fill=CORAL)
    add_text(s, "Objetivo: pipeline GitHub Actions y reporte de drift Evidently",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Pasos",
             Inches(0.5), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Configurar workflow ml.yml con jobs en cadena.",
        "Job 'evaluate' compara modelo nuevo vs Production.",
        "Si supera, mueve a Staging vía API MLflow.",
        "Generar reporte Evidently con dato 'drifted'.",
        "Endpoint /monitor en la API que devuelve resumen.",
    ], Inches(0.7), Inches(2.6), Inches(6), Inches(4), size=14)
    add_text(s, "Comprobaciones",
             Inches(7), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Workflow corre en push y muestra ✓ o ✗.",
        "El job promote solo aplica si f1 nuevo ≥ f1 prod + 1%.",
        "drift.html se genera y se sube como artifact.",
        "GET /monitor responde con drift_score por feature.",
    ], Inches(7.2), Inches(2.6), Inches(6), Inches(4), size=14)

def slide_m5_recap():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Recap Módulo 5", "M5")
    bullets(s, [
        "CI valida código + datos + modelo. CD despliega. CT reentrena.",
        "GitHub Actions / GitLab CI / Argo CD: elige uno y vívelo.",
        "Patrón de despliegue ≠ idea de moda: depende del riesgo y del KPI.",
        "Monitoriza en 4 capas: servicio · datos · modelo · negocio.",
        "Drift en P(X), P(y) o P(y|X). Mide con varios tests, alerta con criterio.",
        "La alerta sin reacción documentada es ruido.",
    ], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5), size=18)

# === MÓDULO 6 — Governance + E2E ===
def slide_feature_store():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Feature Store: el por qué y el cuándo", "M6")
    add_text(s, "Por qué",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        ("Consistencia train/serve", "misma feature, misma definición"),
        ("Reuso entre modelos",      "no recalcular RFM en cada equipo"),
        ("Time-travel correcto",     "evita data leakage temporal"),
        ("Latencia de servicio",     "online store rápido (Redis/DynamoDB)"),
        ("Governance",               "catálogo de features con owner"),
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(4), size=14)

    add_text(s, "Cuándo NO",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=CORAL)
    bullets(s, [
        "1 modelo, 1 equipo, 1 dataset → es overkill.",
        "No hay infra para Redis/Kafka/etc.",
        "El equipo aún no tiene un pipeline maduro.",
    ], Inches(7.2), Inches(1.7), Inches(6), Inches(2.5), size=14)

    add_text(s, "Herramientas",
             Inches(7), Inches(4.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=GREEN)
    bullets(s, [
        ("Feast",     "open source, agnóstico"),
        ("Tecton",    "comercial, gestionado"),
        ("Hopsworks", "open source, completo"),
        ("Vertex FS", "Google Cloud nativo"),
        ("Databricks FS", "integrado en Lakehouse"),
    ], Inches(7.2), Inches(4.7), Inches(6), Inches(2.5), size=13)

def slide_governance():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Model Governance: lo que pide la realidad", "M6")
    items = [
        ("Model Card",          "ficha del modelo: propósito, datos, métricas, sesgos, contraindicaciones"),
        ("Audit trail",         "quién entrenó, quién aprobó, qué dato, qué versión"),
        ("Lineage extremo",     "fila de dato → predicción → decisión → outcome"),
        ("Sesgo y fairness",    "métricas por subgrupo, equal opportunity, demographic parity"),
        ("Explicabilidad",      "shap, LIME, contrafactuales para decisiones de alto impacto"),
        ("Privacidad",          "RGPD, derecho al olvido, anonimización, k-anonymity"),
        ("Human-in-the-loop",   "decisión final humana en uso de alto impacto"),
        ("Reproducción legal",  "se puede recrear el modelo de hace 18 meses si lo pide un juez"),
    ]
    bullets(s, items, Inches(0.5), Inches(1.3), Inches(12.5), Inches(5.5),
            size=14, bullet="✓")

def slide_eu_ai_act():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "EU AI Act + RGPD: lo mínimo que tienes que saber", "M6")
    add_text(s, "EU AI Act (2024)",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "4 niveles de riesgo: inaceptable / alto / limitado / mínimo.",
        "Modelos de alto riesgo (RRHH, salud, crédito, justicia) requieren registro UE.",
        "Obligación de documentación técnica y log de uso.",
        "Transparencia: el usuario debe saber que interactúa con IA.",
    ], Inches(0.7), Inches(1.8), Inches(6), Inches(3), size=13)

    add_text(s, "RGPD aplicado a ML",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "Base legítima para tratar datos de entrenamiento.",
        "Derecho de acceso, rectificación y olvido también afecta a modelos entrenados.",
        "Decisiones automatizadas (Art. 22): derecho a explicación humana.",
        "Privacy by design: anonimización antes que pseudonimización.",
    ], Inches(7.2), Inches(1.8), Inches(6), Inches(3), size=13)

    add_round(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5), fill=CORAL)
    add_text(s, "Buena práctica: cada modelo en producción tiene Model Card + DPIA + log de decisiones",
             Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_finops():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Cost & FinOps en ML", "M6")
    add_text(s, "Dónde se va el dinero",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=CORAL)
    bullets(s, [
        ("Entrenamiento", "GPUs caras, sobre todo si se entrena 'por si acaso'"),
        ("Serving online","instancias siempre encendidas, autoscaling difícil"),
        ("Storage",       "datasets versionados duplicados sin política de retention"),
        ("Egress",        "mover datos entre regiones/providers"),
        ("Tooling SaaS",  "Databricks, W&B, Tecton — facturación opaca"),
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(4), size=13)

    add_text(s, "Palancas para reducir",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=GREEN)
    bullets(s, [
        "Spot instances para entrenamiento.",
        "Batch en lugar de online cuando no hace falta latencia.",
        "Quantización / destilación → modelos más pequeños.",
        "Autoscaling 0→N en horas valle.",
        "Caching agresivo de predicciones repetidas.",
        "FinOps tagging en cada recurso (proyecto, modelo).",
    ], Inches(7.2), Inches(1.7), Inches(6), Inches(4), size=13)

def slide_llmops():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "LLMOps en una slide", "M6")
    cols = [
        ("Prompt",        "Versionar prompts como código (git + tests)"),
        ("Eval",          "Tests de regresión sobre benchmarks propios"),
        ("Guardrails",    "Filtros tóxico/PII, jailbreak, alucinación"),
        ("Coste",         "Tokens in/out, modelo más barato cuando posible"),
        ("RAG",           "Calidad del retrieval > calidad del LLM"),
        ("Drift de uso", "Los usuarios cambian cómo preguntan"),
    ]
    for i, (h, t) in enumerate(cols):
        col = i % 3; row = i // 3
        x = Inches(0.5 + col*4.27); y = Inches(1.4 + row*2.7)
        add_round(s, x, y, Inches(4.1), Inches(2.4), fill=PURPLE)
        add_rect(s, x, y, Inches(4.1), Inches(0.7), fill=NAVY)
        add_text(s, h, x, y, Inches(4.1), Inches(0.7),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, t, x+Inches(0.3), y+Inches(0.9), Inches(3.5), Inches(1.5),
                 size=14, color=WHITE)

def slide_lab5_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Lab 5 — Recorrido end-to-end (45 min)", "M6")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7), fill=GREEN)
    add_text(s, "Objetivo: ejecutar el ciclo completo y discutir qué falta",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "El recorrido",
             Inches(0.5), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "git pull → dvc pull → dependencias listas.",
        "pytest verde.",
        "python src/train.py loggea run en MLflow.",
        "promote_if_better.py mueve a Staging.",
        "docker compose up: API contesta /predict.",
        "evidently report sobre datos sintéticos drifted.",
    ], Inches(0.7), Inches(2.6), Inches(6), Inches(4), size=14)

    add_text(s, "Discusión guiada",
             Inches(7), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "¿Qué falta para producción 'real'?",
        "Secret management.",
        "Observabilidad central (Grafana/Prom).",
        "Rollback automatizado.",
        "Pruebas de carga periódicas.",
        "Disaster recovery del Registry.",
    ], Inches(7.2), Inches(2.6), Inches(6), Inches(4), size=14)

def slide_roadmap_learning():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Hoja de ruta de aprendizaje", "fin")
    levels = [
        ("Junior",    "Sabe usar MLflow, DVC y FastAPI. Entrega un modelo dockerizado.", CYAN),
        ("Mid",       "Diseña pipelines CI/CD y monitorea drift. Domina al menos un orquestador.", CORAL),
        ("Senior",    "Diseña la plataforma MLOps. Decide trade-offs de arquitectura, coste y riesgo.", GREEN),
        ("Staff/Arq.","Define estándares de governance y forma a otros equipos. Habla con legal y producto.", NAVY),
    ]
    for i, (n, d, c) in enumerate(levels):
        y = Inches(1.4 + i*1.35)
        add_round(s, Inches(0.5), y, Inches(2.3), Inches(1.15), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2.3), Inches(1.15),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(2.9), y, Inches(9.9), Inches(1.15), fill=WHITE, line=c)
        add_text(s, d, Inches(3.1), y, Inches(9.6), Inches(1.15),
                 size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)

def slide_resources():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Recursos para seguir", "fin")
    add_text(s, "Libros",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "Designing ML Systems — Chip Huyen (O'Reilly).",
        "Practical MLOps — Noah Gift (O'Reilly).",
        "Reliable Machine Learning — Cathy Chen (O'Reilly).",
        "Building ML Powered Applications — Emmanuel Ameisen.",
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(3), size=14)
    add_text(s, "Online",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "ml-ops.org (patrones y referencias).",
        "Made With ML (curso open source).",
        "MLOps Community (Slack + podcast).",
        "Google MLOps whitepaper.",
        "Andrej Karpathy 'Software 2.0'.",
    ], Inches(7.2), Inches(1.7), Inches(6), Inches(4), size=14)
    add_text(s, "Comunidades en español",
             Inches(0.5), Inches(5.0), Inches(12), Inches(0.5),
             size=20, bold=True, color=CORAL)
    bullets(s, [
        "ANBAN (claro!).",
        "BigData Spain, Codemotion, Spain AI.",
        "PyData Spain, Madrid Pandas, Barcelona ML Meetup.",
    ], Inches(0.7), Inches(5.5), Inches(12), Inches(2), size=14)

def slide_thanks():
    s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.6), SH, fill=CORAL)
    add_text(s, "Gracias",
             Inches(1), Inches(2.5), Inches(11), Inches(1.5),
             size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Preguntas, comentarios y trolling bienvenidos",
             Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             size=22, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(s, "José Picón · jose.bobal@gmail.com",
             Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Material del curso disponible en el repo de ANBAN",
             Inches(1), Inches(6.2), Inches(11), Inches(0.6),
             size=14, color=GREY, align=PP_ALIGN.CENTER)

# Lista de slides Día 2
slides_funcs = [
    slide_cover,                                     # 1
    slide_d1_recap,                                  # 2
    slide_agenda_d2,                                 # 3
    lambda: slide_module_divider(4, "Empaquetado y serving",
        "Llevar el modelo del Registry a una API en contenedor", CYAN),  # 4
    slide_serving_patterns,                          # 5
    slide_packaging_options,                         # 6
    slide_api_design,                                # 7
    slide_perf_metrics,                              # 8
    slide_train_serve_skew,                          # 9
    slide_lab3_intro,                                # 10
    slide_dockerfile,                                # 11
    slide_m4_recap,                                  # 12
    lambda: slide_module_divider(5, "CI/CD + Monitoring + Drift",
        "Automatizar la release y detectar la degradación", CORAL),  # 13
    slide_ci_cd_ct,                                  # 14
    slide_pipeline_jobs,                             # 15
    slide_github_actions,                            # 16
    slide_deploy_strategies,                         # 17
    slide_monitoring_what,                           # 18
    slide_drift_types,                               # 19
    slide_drift_tests,                               # 20
    slide_evidently,                                 # 21
    slide_alerting,                                  # 22
    slide_lab4_intro,                                # 23
    slide_m5_recap,                                  # 24
    lambda: slide_module_divider(6, "Caso end-to-end + Governance",
        "Cerrar el ciclo y mirar lo que falta", GREEN),  # 25
    slide_feature_store,                             # 26
    slide_governance,                                # 27
    slide_eu_ai_act,                                 # 28
    slide_finops,                                    # 29
    slide_llmops,                                    # 30
    slide_lab5_intro,                                # 31
    slide_roadmap_learning,                          # 32
    slide_resources,                                 # 33
    slide_thanks,                                    # 34
]

for f in slides_funcs:
    f()

total = len(prs.slides)
for i, slide in enumerate(prs.slides):
    if i == 0:
        continue
    footer_pagenum(slide, i+1, total)

out = "/home/jmpicon/Documentos/Anbam/slides/dia2_mlops_dataops.pptx"
prs.save(out)
print(f"OK · {total} slides · {out}")
