"""
Generador del PPTX Día 1 — MLOps/DataOps (ANBAN).
Genera slides/dia1_mlops_dataops.pptx con tema profesional.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# === PALETA ANBAN ===
NAVY = RGBColor(0x0B, 0x1F, 0x3A)       # Fondo títulos
CYAN = RGBColor(0x00, 0xB4, 0xD8)       # Acento
CORAL = RGBColor(0xFF, 0x6B, 0x35)      # Highlight
BG = RGBColor(0xF8, 0xFA, 0xFC)         # Fondo claro
INK = RGBColor(0x1A, 0x1D, 0x29)        # Texto
GREY = RGBColor(0x6B, 0x72, 0x80)       # Texto secundario
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_text(slide, text, left, top, width, height, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_rect(slide, left, top, width, height, *, fill=NAVY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s

def add_round(slide, left, top, width, height, *, fill=CYAN, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s

def header(slide, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), fill=NAVY)
    # franja cyan
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.06), fill=CYAN)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(11), Inches(0.6),
             size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        add_text(slide, kicker, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.45),
                 size=12, color=CYAN, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, bold=True)
    # footer
    add_text(slide, "MLOps / DataOps · ANBAN · José Picón",
             Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             size=10, color=GREY)

def footer_pagenum(slide, n, total):
    add_text(slide, f"{n} / {total}",
             Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             size=10, color=GREY, align=PP_ALIGN.RIGHT)

def bullets(slide, items, left, top, width, height, *, size=18, color=INK, bullet="●"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        if isinstance(item, tuple):
            head, sub = item
            r = p.add_run(); r.text = f"{bullet}  "
            r.font.size = Pt(size); r.font.color.rgb = CYAN; r.font.bold = True
            r2 = p.add_run(); r2.text = head
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.bold = True
            r3 = p.add_run(); r3.text = f" — {sub}"
            r3.font.size = Pt(size-2); r3.font.color.rgb = GREY
        else:
            r = p.add_run(); r.text = f"{bullet}  "
            r.font.size = Pt(size); r.font.color.rgb = CYAN; r.font.bold = True
            r2 = p.add_run(); r2.text = item
            r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

# -----------------------------------------------------------------------------
# SLIDES
# -----------------------------------------------------------------------------
slides_data = []

def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    # banda izquierda decorativa
    add_rect(s, 0, 0, Inches(0.6), SH, fill=CYAN)
    # título grande
    add_text(s, "MLOps / DataOps", Inches(1.0), Inches(2.0), Inches(11), Inches(1.4),
             size=60, bold=True, color=WHITE)
    add_text(s, "De la prueba de concepto a producción", Inches(1.0), Inches(3.2), Inches(11), Inches(0.8),
             size=28, color=CYAN)
    # línea
    add_rect(s, Inches(1.0), Inches(4.1), Inches(3), Inches(0.05), fill=CORAL)
    add_text(s, "Día 1 — DataOps y fundamentos de MLOps", Inches(1.0), Inches(4.3), Inches(11), Inches(0.6),
             size=22, color=WHITE)
    add_text(s, "ANBAN · Asociación Big Data España", Inches(1.0), Inches(6.0), Inches(11), Inches(0.5),
             size=16, color=CYAN)
    add_text(s, "José Picón · 2026", Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             size=14, color=WHITE)

def slide_about_speaker():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Quién soy y qué vamos a hacer", "intro")
    add_text(s, "Sobre el docente",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=22, bold=True, color=NAVY)
    bullets(s, [
        "Ingeniero de seguridad y datos.",
        "Docente Curso Especialización ENS y formaciones de Big Data.",
        "Experiencia llevando modelos a producción en sectores regulados.",
    ], Inches(0.5), Inches(1.8), Inches(6), Inches(2.5), size=16)

    add_text(s, "Cómo trabajaremos",
             Inches(7.0), Inches(1.2), Inches(6), Inches(0.5),
             size=22, bold=True, color=NAVY)
    bullets(s, [
        "Teoría densa pero al grano (≈30%).",
        "Laboratorios reproducibles con Docker (≈60%).",
        "Discusión y dudas en cualquier momento (≈10%).",
        "Todo el material en el repositorio del curso.",
    ], Inches(7.0), Inches(1.8), Inches(6), Inches(3), size=16)

    add_round(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.0), fill=CORAL)
    add_text(s, "Regla del curso: ningún modelo en producción si no podemos contar su trazabilidad",
             Inches(0.7), Inches(5.6), Inches(12), Inches(1.0),
             size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

def slide_agenda_d1():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Agenda · Día 1 (4 horas)", "agenda")
    rows = [
        ("00:00–00:45", "M1", "Introducción MLOps/DataOps y madurez", CYAN),
        ("00:45–01:00", "—",  "Pausa", GREY),
        ("01:00–02:15", "M2", "DataOps: pipelines, DVC, Great Expectations", CORAL),
        ("02:15–02:30", "—",  "Pausa", GREY),
        ("02:30–03:45", "M3", "Experiment tracking con MLflow", GREEN),
        ("03:45–04:00", "—",  "Cierre y reto del día", NAVY),
    ]
    y = Inches(1.4)
    for h, code, txt, col in rows:
        add_round(s, Inches(0.5), y, Inches(1.7), Inches(0.7), fill=col)
        add_text(s, h, Inches(0.5), y, Inches(1.7), Inches(0.7),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(2.4), y, Inches(0.9), Inches(0.7), fill=NAVY)
        add_text(s, code, Inches(2.4), y, Inches(0.9), Inches(0.7),
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, txt, Inches(3.5), y, Inches(9.5), Inches(0.7),
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.85)

def slide_objectives():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Objetivos del curso", "intro")
    objs = [
        "Diagnosticar la madurez MLOps de tu organización (L0/L1/L2).",
        "Versionar datos y validarlos automáticamente (DVC + GE).",
        "Trazar y reproducir cualquier experimento (MLflow + Registry).",
        "Servir un modelo como API en contenedor (FastAPI + Docker).",
        "Automatizar el pipeline ML completo con CI/CD (GitHub Actions).",
        "Detectar drift y degradación (Evidently) y reaccionar.",
        "Aplicar governance: feature store, model cards, RGPD/AI Act.",
    ]
    add_text(s, "Al terminar el curso serás capaz de…",
             Inches(0.5), Inches(1.2), Inches(12), Inches(0.6),
             size=20, bold=True, color=NAVY)
    bullets(s, objs, Inches(0.7), Inches(2.0), Inches(12), Inches(5), size=18)

# ============================================================================
# MÓDULO 1 — INTRODUCCIÓN
# ============================================================================
def slide_module_divider(num, title, subtitle, color):
    s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
    add_rect(s, 0, Inches(2.5), SW, Inches(2.5), fill=color)
    add_text(s, f"Módulo {num}", Inches(0.8), Inches(2.7), Inches(12), Inches(0.6),
             size=22, color=WHITE, bold=True)
    add_text(s, title, Inches(0.8), Inches(3.2), Inches(12), Inches(1.1),
             size=44, color=WHITE, bold=True)
    add_text(s, subtitle, Inches(0.8), Inches(5.2), Inches(12), Inches(0.6),
             size=20, color=CYAN)

def slide_why_mlops():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "¿Por qué existe MLOps?", "M1")
    add_text(s, "El 87% de los proyectos de ML nunca llegan a producción",
             Inches(0.5), Inches(1.1), Inches(12), Inches(0.6),
             size=22, bold=True, color=CORAL)
    add_text(s, "(VentureBeat / Algorithmia · 2020)",
             Inches(0.5), Inches(1.6), Inches(12), Inches(0.4),
             size=12, color=GREY)
    # 4 cajas
    boxes = [
        ("Datos", "Calidad inestable, sin versionado, sin lineage"),
        ("Reproducibilidad", "Notebooks irrepetibles, semillas no fijadas"),
        ("Despliegue", "Train ≠ serve, sin contrato de API"),
        ("Operación", "Sin monitoring, sin drift, sin rollback"),
    ]
    for i, (h, t) in enumerate(boxes):
        x = Inches(0.5 + i*3.13); y = Inches(2.5)
        add_round(s, x, y, Inches(2.95), Inches(2.0), fill=WHITE, line=CYAN)
        add_rect(s, x, y, Inches(2.95), Inches(0.5), fill=NAVY)
        add_text(s, h, x, y, Inches(2.95), Inches(0.5),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, t, x+Inches(0.15), y+Inches(0.7), Inches(2.7), Inches(1.3),
                 size=14, color=INK)
    add_round(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), fill=BG, line=CORAL)
    add_text(s, "MLOps = la disciplina que se asegura de que el modelo llegue a producción y se mantenga útil",
             Inches(0.7), Inches(5.0), Inches(12), Inches(1.5),
             size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

def slide_hidden_debt():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Deuda técnica oculta en sistemas ML", "M1")
    add_text(s, 'Sculley et al., NeurIPS 2015 · "Hidden Technical Debt in ML Systems"',
             Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
             size=14, color=GREY)
    # diagrama tipo "ML code es la cajita pequeña"
    add_round(s, Inches(0.5), Inches(2.0), Inches(7.8), Inches(4.5), fill=WHITE, line=NAVY)
    # cajitas internas
    deps = [
        ("Configuración", 0.7, 2.2, 1.8, 0.6),
        ("Recolección de datos", 2.7, 2.2, 2.4, 0.6),
        ("Análisis de datos", 5.3, 2.2, 2.0, 0.6),
        ("Recursos máquina", 0.7, 3.0, 2.5, 0.6),
        ("Tooling serving", 3.4, 3.0, 2.0, 0.6),
        ("Monitoring", 5.6, 3.0, 1.7, 0.6),
        ("Pipelines de datos", 0.7, 3.8, 2.5, 0.6),
        ("Tests automatizados", 3.4, 3.8, 2.5, 0.6),
        ("Gestión de procesos", 0.7, 4.6, 3.0, 0.6),
        ("Análisis de errores", 4.0, 4.6, 3.2, 0.6),
        ("Verificación de features", 0.7, 5.4, 3.5, 0.6),
        ("Sistemas de despliegue", 4.4, 5.4, 2.8, 0.6),
    ]
    for txt, x, y, w, h in deps:
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=BG, line=GREY)
        add_text(s, txt, Inches(x), Inches(y), Inches(w), Inches(h),
                 size=11, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # ML code como caja roja pequeña
    add_rect(s, Inches(3.7), Inches(3.85), Inches(1.4), Inches(0.5), fill=CORAL)
    add_text(s, "ML Code", Inches(3.7), Inches(3.85), Inches(1.4), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Lectura",
             Inches(8.6), Inches(2.0), Inches(4.5), Inches(0.4),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "El código del modelo es una fracción del sistema.",
        "El resto es infraestructura, datos y operación.",
        "MLOps gestiona toda la caja, no solo el código.",
        "Ignorar lo que rodea al modelo es la causa #1 de fallo.",
    ], Inches(8.6), Inches(2.5), Inches(4.5), Inches(4), size=15)

def slide_three_dimensions():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Las tres dimensiones de MLOps", "M1")
    cols = [
        ("CÓDIGO", "Software tradicional. Git, tests, build.", "DevOps", CYAN),
        ("DATOS", "Cambian, crecen, se contaminan. Versionado y validación.", "DataOps", CORAL),
        ("MODELO", "Función entrenada del par (datos, código). Reentrena.", "ModelOps", GREEN),
    ]
    for i, (h, sub, role, col) in enumerate(cols):
        x = Inches(0.5 + i*4.27)
        add_round(s, x, Inches(1.4), Inches(4.1), Inches(4.5), fill=WHITE, line=col)
        add_rect(s, x, Inches(1.4), Inches(4.1), Inches(0.7), fill=col)
        add_text(s, h, x, Inches(1.4), Inches(4.1), Inches(0.7),
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sub, x+Inches(0.2), Inches(2.3), Inches(3.7), Inches(2),
                 size=15, color=INK)
        add_text(s, f"Disciplina: {role}", x+Inches(0.2), Inches(5.2), Inches(3.7), Inches(0.5),
                 size=13, bold=True, color=col)
    add_round(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), fill=NAVY)
    add_text(s, "MLOps = unión coherente de las tres dimensiones",
             Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_devops_vs():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "DevOps vs DataOps vs MLOps vs LLMOps", "M1")
    # Tabla simple
    rows = [
        ("",           "DevOps",          "DataOps",            "MLOps",                  "LLMOps"),
        ("Artefacto",  "Aplicación",      "Pipeline de datos",  "Modelo + datos + código","LLM + prompt + RAG"),
        ("Cambio",     "Determinista",    "Datos cambian",      "Datos y modelo cambian", "Modelo opaco, prompt cambia"),
        ("CI/CD",      "Test, build, deploy","Test datos, dq",  "+ Train + Eval + Promote","+ Eval LLM, regresión"),
        ("Monitoring", "Latencia, errores","Calidad, freshness","+ Drift, performance",   "+ Tóxico, alucinación, coste"),
        ("Versionado", "Git",             "Git + DVC",          "Git + DVC + Registry",   "+ prompt registry"),
    ]
    rh = Inches(0.65)
    cw = [Inches(2.0), Inches(2.5), Inches(2.7), Inches(3.0), Inches(3.0)]
    x0 = Inches(0.4); y0 = Inches(1.3)
    x = x0
    for ci, c in enumerate(rows[0]):
        fill = NAVY if ci > 0 else CYAN
        add_rect(s, x, y0, cw[ci], rh, fill=fill)
        add_text(s, c, x, y0, cw[ci], rh, size=14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cw[ci]
    for ri, row in enumerate(rows[1:]):
        y = y0 + rh*(ri+1)
        x = x0
        for ci, c in enumerate(row):
            fill = BG if ri % 2 == 0 else WHITE
            add_rect(s, x, y, cw[ci], rh, fill=fill, line=GREY)
            add_text(s, c, x+Inches(0.1), y, cw[ci], rh, size=12,
                     color=INK, anchor=MSO_ANCHOR.MIDDLE,
                     bold=(ci == 0))
            x += cw[ci]

def slide_crisp_dm():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "CRISP-DM → CRISP-ML(Q): cambia el ciclo", "M1")
    # circulo a la izquierda
    add_text(s, "CRISP-DM (1996)",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=GREY)
    bullets(s, [
        "Business Understanding",
        "Data Understanding",
        "Data Preparation",
        "Modeling",
        "Evaluation",
        "Deployment",
    ], Inches(0.7), Inches(1.8), Inches(6), Inches(4), size=16)
    add_text(s, "Pensado para entregables únicos.\nNo contempla mantener el modelo.",
             Inches(0.5), Inches(5.8), Inches(6), Inches(1.2),
             size=13, color=GREY)

    add_text(s, "CRISP-ML(Q) (2021)",
             Inches(7.0), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        ("Business & Data", "objetivos y datos viables"),
        ("Data Engineering", "data quality requirements"),
        ("ML Modeling", "+ quality gates"),
        ("Evaluation", "robustez, fairness, explicabilidad"),
        ("Deployment", "patrones de release"),
        ("Monitoring & Maintenance", "el bucle, no el final"),
    ], Inches(7.2), Inches(1.8), Inches(6), Inches(4.8), size=15)
    add_round(s, Inches(7.0), Inches(6.4), Inches(5.8), Inches(0.6), fill=CORAL)
    add_text(s, "Calidad como requisito explícito en cada fase",
             Inches(7.0), Inches(6.4), Inches(5.8), Inches(0.6),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_maturity(level, title, characteristics, problems, color):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, f"Madurez MLOps · {title}", "M1")
    add_round(s, Inches(0.5), Inches(1.2), Inches(2.5), Inches(2.5), fill=color)
    add_text(s, level, Inches(0.5), Inches(1.2), Inches(2.5), Inches(2.5),
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, Inches(3.2), Inches(1.4), Inches(9), Inches(0.6),
             size=24, bold=True, color=NAVY)
    add_text(s, "Características", Inches(3.2), Inches(2.0), Inches(9), Inches(0.4),
             size=16, bold=True, color=color)
    bullets(s, characteristics, Inches(3.4), Inches(2.5), Inches(9), Inches(2.5), size=15)
    add_text(s, "Síntomas/Problemas",
             Inches(0.5), Inches(4.2), Inches(12), Inches(0.4),
             size=18, bold=True, color=CORAL)
    bullets(s, problems, Inches(0.7), Inches(4.7), Inches(12), Inches(2.5),
            size=15, color=INK, bullet="⚠")

def slide_roles():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Roles y responsabilidades", "M1")
    roles = [
        ("Data Engineer", "Pipelines, ingesta, calidad, DataOps", CYAN),
        ("Data Scientist", "Hipótesis, EDA, modelado, métrica negocio", CORAL),
        ("ML Engineer", "Pone el modelo en producción, optimiza", GREEN),
        ("MLOps Engineer", "Plataforma: tracking, registry, CI/CD, infra ML", NAVY),
        ("SRE / Platform", "Infra cluster, observabilidad, costes", RGBColor(0xA0, 0x55, 0xF7)),
        ("Product / Negocio", "Caso de uso, KPI, prioridades", RGBColor(0xF5, 0x9E, 0x0B)),
    ]
    for i, (r, d, c) in enumerate(roles):
        col = i % 3; row = i // 3
        x = Inches(0.5 + col*4.27); y = Inches(1.4 + row*2.7)
        add_round(s, x, y, Inches(4.1), Inches(2.4), fill=WHITE, line=c)
        add_rect(s, x, y, Inches(4.1), Inches(0.6), fill=c)
        add_text(s, r, x, y, Inches(4.1), Inches(0.6),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, d, x+Inches(0.2), y+Inches(0.8), Inches(3.7), Inches(1.5),
                 size=14, color=INK)

def slide_modern_stack():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Stack moderno típico (open source)", "M1")
    layers = [
        ("Datos",          ["Airbyte", "dbt", "DVC", "LakeFS", "Delta"], CORAL),
        ("Calidad",        ["Great Expectations", "Soda", "Pandera"], CYAN),
        ("Orquestación",   ["Airflow", "Prefect", "Dagster", "Argo"], GREEN),
        ("Experiment",     ["MLflow", "Weights & Biases", "Neptune"], NAVY),
        ("Feature Store",  ["Feast", "Tecton", "Hopsworks"], RGBColor(0xA0, 0x55, 0xF7)),
        ("Serving",        ["FastAPI", "BentoML", "TF Serving", "TorchServe"], RGBColor(0xF5, 0x9E, 0x0B)),
        ("Monitoring",     ["Evidently", "WhyLabs", "Arize", "NannyML"], RGBColor(0xEC, 0x49, 0x99)),
        ("CI/CD",          ["GitHub Actions", "GitLab CI", "Jenkins", "ArgoCD"], RGBColor(0x06, 0xB6, 0xD4)),
    ]
    for i, (layer, tools, c) in enumerate(layers):
        y = Inches(1.3 + i*0.7)
        add_rect(s, Inches(0.5), y, Inches(2.5), Inches(0.6), fill=c)
        add_text(s, layer, Inches(0.5), y, Inches(2.5), Inches(0.6),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.0), y, Inches(9.8), Inches(0.6), fill=BG, line=GREY)
        add_text(s, " · ".join(tools), Inches(3.1), y, Inches(9.7), Inches(0.6),
                 size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)

def slide_diagnostic_exercise():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Ejercicio en clase: ¿en qué nivel está tu organización?", "M1")
    add_text(s, "Marca con un check (en grupo, 8 minutos)",
             Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             size=16, color=GREY, bold=True)
    items = [
        "Existe un pipeline reproducible de datos (no hay 'el script de Juan').",
        "El dataset usado para entrenar está versionado y se puede recuperar.",
        "Cada experimento queda registrado: params, métricas y modelo.",
        "Hay un modelo en producción que NO sea un .pkl en el laptop de alguien.",
        "El modelo se sirve detrás de una API con contrato versionado.",
        "Existe un test automático que se ejecuta antes de promocionar un modelo.",
        "Se vigila la distribución de los inputs en producción (drift).",
        "Hay procedimiento de rollback de modelo documentado y probado.",
        "Cada predicción se puede trazar a (modelo, dataset, código).",
        "Existe un Model Registry consultable por alguien que no eres tú.",
    ]
    bullets(s, items, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
            size=14, bullet="☐")
    add_round(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6), fill=CORAL)
    add_text(s, "0–3 → L0  ·  4–6 → L1  ·  7–10 → L2",
             Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# MÓDULO 2 — DATAOPS
# ============================================================================
def slide_dataops_manifesto():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "DataOps: principios fundacionales", "M2")
    add_text(s, "DataOps = DevOps + Agile + Lean Manufacturing aplicados al dato",
             Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
             size=18, bold=True, color=NAVY)
    principles = [
        ("Continuous integration de datos", "Cada cambio en el dato se prueba"),
        ("Observabilidad", "Métricas, logs y lineage del dato"),
        ("Reproducibilidad", "Mismo input → mismo output"),
        ("Automatización", "Cero pasos manuales en el camino crítico"),
        ("Calidad por diseño", "Tests de datos antes de tests de modelo"),
        ("Iteración rápida", "Pequeños cambios validados, no big bang"),
    ]
    for i, (h, t) in enumerate(principles):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(2.0 + row*1.55)
        add_round(s, x, y, Inches(6.2), Inches(1.4), fill=WHITE, line=CYAN)
        add_rect(s, x, y, Inches(0.15), Inches(1.4), fill=CYAN)
        add_text(s, h, x+Inches(0.3), y+Inches(0.1), Inches(5.8), Inches(0.4),
                 size=16, bold=True, color=NAVY)
        add_text(s, t, x+Inches(0.3), y+Inches(0.6), Inches(5.8), Inches(0.7),
                 size=13, color=GREY)

def slide_data_pipeline():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Anatomía de un pipeline de datos", "M2")
    stages = [("Ingesta", CYAN), ("Validación", CORAL), ("Limpieza", GREEN),
              ("Feature Eng.", NAVY), ("Almacén", RGBColor(0xA0,0x55,0xF7)),
              ("Servir", RGBColor(0xF5,0x9E,0x0B))]
    x = Inches(0.4)
    for i, (name, col) in enumerate(stages):
        add_round(s, x, Inches(2.0), Inches(1.9), Inches(1.3), fill=col)
        add_text(s, name, x, Inches(2.0), Inches(1.9), Inches(1.3),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(stages)-1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     x+Inches(1.95), Inches(2.4), Inches(0.18), Inches(0.5))
            arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background()
        x += Inches(2.13)
    # debajo, qué se valida en cada stage
    detail = [
        ("Ingesta",      "Schema, conteos vs origen, duplicados"),
        ("Validación",   "Reglas de negocio, rangos, dominios"),
        ("Limpieza",     "Nulos, outliers, normalización"),
        ("Feature Eng.", "Lineage de cada feature, no leakage"),
        ("Almacén",      "Particionado, retention, formato (parquet)"),
        ("Servir",       "Train ≠ serve es un bug, no una feature"),
    ]
    for i, (h, t) in enumerate(detail):
        col = i % 3; row = i // 3
        x = Inches(0.5 + col*4.27); y = Inches(3.7 + row*1.5)
        add_round(s, x, y, Inches(4.1), Inches(1.3), fill=BG, line=GREY)
        add_text(s, h, x+Inches(0.2), y+Inches(0.1), Inches(3.9), Inches(0.4),
                 size=15, bold=True, color=NAVY)
        add_text(s, t, x+Inches(0.2), y+Inches(0.55), Inches(3.9), Inches(0.7),
                 size=13, color=INK)

def slide_why_data_versioning():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "¿Por qué versionar datos?", "M2")
    bullets(s, [
        "Reproducir un experimento sin la versión exacta del dataset = imposible.",
        "Auditar una decisión del modelo = imposible si no sabes qué dato lo entrenó.",
        "Volver atrás cuando una nueva muestra contamina el modelo.",
        "Comparar modelos: solo es justo si comparten la misma versión de datos.",
        "Cumplir RGPD / AI Act: derecho al olvido implica saber qué datos se usaron.",
    ], Inches(0.5), Inches(1.2), Inches(12), Inches(3), size=17)
    add_round(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.5), fill=NAVY)
    add_text(s, "Pero Git no sirve para datasets", Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
             size=18, bold=True, color=CORAL)
    bullets(s, [
        "Git almacena cada versión completa: pesos, tarda, y revienta el repo.",
        "Diffs binarios: inútiles.",
        "Repos con >1GB son una pesadilla de clonado.",
    ], Inches(0.7), Inches(5.1), Inches(12), Inches(2),
            size=14, color=WHITE)

def slide_dvc_alternatives():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Versionado de datos: DVC y alternativas", "M2")
    tools = [
        ("DVC",       "Versiona ficheros con punteros en Git + storage (S3, GCS, local)",
         "ML clásico, datasets de archivos", CYAN),
        ("LakeFS",    "Branches y commits sobre object storage (al estilo Git)",
         "Data lakes en S3 con flujos de PR", CORAL),
        ("Delta Lake","Tablas con time-travel y ACID sobre Parquet",
         "Data warehouse moderno (Spark, Databricks)", GREEN),
        ("lakeFS + DuckDB", "Combinación ligera para sandbox local",
         "Equipos pequeños sin Spark", NAVY),
    ]
    for i, (n, d, use, c) in enumerate(tools):
        y = Inches(1.3 + i*1.4)
        add_round(s, Inches(0.5), y, Inches(2.5), Inches(1.2), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2.5), Inches(1.2),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(3.1), y, Inches(9.7), Inches(1.2), fill=WHITE, line=c)
        add_text(s, d, Inches(3.3), y+Inches(0.1), Inches(9.4), Inches(0.5),
                 size=14, color=INK)
        add_text(s, f"Cuándo: {use}", Inches(3.3), y+Inches(0.65), Inches(9.4), Inches(0.5),
                 size=13, color=c, bold=True)

def slide_dvc_internals():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Cómo funciona DVC bajo el capó", "M2")
    bullets(s, [
        ("dvc add data.csv", "calcula hash MD5, mueve el fichero a la cache (.dvc/cache)"),
        ("genera data.csv.dvc", "fichero pequeño con el hash → versionable en Git"),
        ("dvc remote", "configura S3/GCS/MinIO como almacenamiento real"),
        ("dvc push / pull", "sincroniza los blobs grandes con el remote"),
        ("dvc.yaml", "define stages como un Makefile reproducible (DAG)"),
        ("dvc repro", "vuelve a ejecutar solo lo que cambió"),
    ], Inches(0.5), Inches(1.2), Inches(12), Inches(4), size=15)

    add_round(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2), fill=NAVY)
    code = ("# dvc.yaml\n"
            "stages:\n"
            "  validate:\n"
            "    cmd: python validate.py raw/data.csv\n"
            "    deps: [raw/data.csv, validate.py]\n"
            "    outs: [validated/data.parquet]")
    add_text(s, code, Inches(0.7), Inches(5.1), Inches(12), Inches(1.8),
             size=14, color=CYAN, font="Consolas")

def slide_data_quality_dimensions():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Calidad de datos: dimensiones a medir", "M2")
    dims = [
        ("Completitud",  "% de campos no nulos sobre los esperados"),
        ("Validez",      "Cumple el schema, dominio y reglas de negocio"),
        ("Unicidad",     "No duplicados en claves naturales"),
        ("Consistencia", "Coherente entre fuentes (ID en CRM = ID en facturación)"),
        ("Precisión",    "Refleja la realidad (ej. coordenadas correctas)"),
        ("Frescura",     "Latencia desde el evento hasta que está disponible"),
        ("Trazabilidad", "Lineage: ¿de dónde viene este registro?"),
    ]
    for i, (h, t) in enumerate(dims):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(1.3 + row*0.85)
        add_round(s, x, y, Inches(6.2), Inches(0.75), fill=WHITE, line=CYAN)
        add_rect(s, x, y, Inches(2), Inches(0.75), fill=CYAN)
        add_text(s, h, x, y, Inches(2), Inches(0.75),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, t, x+Inches(2.1), y, Inches(4.0), Inches(0.75),
                 size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)

def slide_great_expectations():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Great Expectations: tests para datos", "M2")
    add_text(s, "Concepto", Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "Suite de expectations declarativas (assertion sobre el dato).",
        "Genera Data Docs HTML como informe legible.",
        "Se integra con Airflow, Prefect, dbt y Spark.",
        "Validation Result auditable en cada ejecución.",
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(3), size=14)

    add_text(s, "Ejemplo", Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    add_round(s, Inches(7), Inches(1.7), Inches(6), Inches(5), fill=NAVY)
    code = ("import great_expectations as gx\n"
            "ctx = gx.get_context()\n"
            "df = ctx.data_sources['files'].read_csv('data.csv')\n"
            "\n"
            "df.expect_column_values_to_not_be_null('user_id')\n"
            "df.expect_column_values_to_be_unique('user_id')\n"
            "df.expect_column_values_to_be_between(\n"
            "    'age', min_value=0, max_value=120)\n"
            "df.expect_column_values_to_match_regex(\n"
            "    'email', r'^[\\w.-]+@[\\w.-]+$')\n"
            "\n"
            "result = df.validate()\n"
            "assert result.success")
    add_text(s, code, Inches(7.2), Inches(1.85), Inches(5.7), Inches(4.7),
             size=12, color=CYAN, font="Consolas")

def slide_pandera_soda():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Alternativas: Pandera y Soda", "M2")
    add_text(s, "Pandera", Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=22, bold=True, color=NAVY)
    add_text(s, "Validación tipada para DataFrames",
             Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
             size=13, color=GREY)
    add_round(s, Inches(0.5), Inches(2.2), Inches(6), Inches(3.3), fill=NAVY)
    code1 = ("import pandera as pa\n\n"
             "schema = pa.DataFrameSchema({\n"
             '    "age":    pa.Column(int, pa.Check.in_range(0, 120)),\n'
             '    "email":  pa.Column(str, pa.Check.str_matches(r"@")),\n'
             '    "amount": pa.Column(float, pa.Check.gt(0)),\n'
             "})\n\n"
             "schema.validate(df)")
    add_text(s, code1, Inches(0.7), Inches(2.3), Inches(5.6), Inches(3.1),
             size=12, color=CYAN, font="Consolas")
    bullets(s, [
        "Sintaxis pythónica, integra con FastAPI",
        "Bueno cuando el equipo ya escribe pandas",
    ], Inches(0.7), Inches(5.6), Inches(6), Inches(1.5), size=13)

    add_text(s, "Soda Core", Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=22, bold=True, color=NAVY)
    add_text(s, "DSL declarativo en YAML, ideal para SQL warehouses",
             Inches(7), Inches(1.7), Inches(6), Inches(0.4),
             size=13, color=GREY)
    add_round(s, Inches(7), Inches(2.2), Inches(6), Inches(3.3), fill=NAVY)
    code2 = ("# checks.yml\n"
             "checks for orders:\n"
             "  - row_count > 0\n"
             "  - missing_count(customer_id) = 0\n"
             "  - duplicate_count(order_id) = 0\n"
             "  - freshness(created_at) < 1d\n"
             "  - avg(amount) between 10 and 10000")
    add_text(s, code2, Inches(7.2), Inches(2.3), Inches(5.6), Inches(3.1),
             size=12, color=CYAN, font="Consolas")
    bullets(s, [
        "Cli y servicio (Soda Cloud)",
        "Pensado para warehouses (Snowflake, BigQuery, Postgres)",
    ], Inches(7.2), Inches(5.6), Inches(6), Inches(1.5), size=13)

def slide_orchestrators():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Orquestadores: ¿cuál uso?", "M2")
    rows = [
        ("Airflow",  "Veterano, gran ecosistema, UI clásica",
         "ETL pesado, equipos grandes, ya en uso", CYAN),
        ("Prefect",  "Pythonic, dynamic DAGs, observabilidad fina",
         "Equipos data engineers que viven en Python", CORAL),
        ("Dagster",  "Asset-oriented, SDA, types, integra dbt",
         "Equipos modernos que piensan en assets, no en tasks", GREEN),
        ("Argo Wf",  "Kubernetes nativo, contenedores como tasks",
         "Cargas ML pesadas en K8s, ya hay plataforma K8s", NAVY),
    ]
    for i, (n, d, u, c) in enumerate(rows):
        y = Inches(1.3 + i*1.4)
        add_round(s, Inches(0.5), y, Inches(2), Inches(1.2), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2), Inches(1.2),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(2.6), y, Inches(10.2), Inches(1.2), fill=WHITE, line=c)
        add_text(s, d, Inches(2.8), y+Inches(0.1), Inches(10), Inches(0.5),
                 size=14, color=INK)
        add_text(s, f"Cuándo: {u}", Inches(2.8), y+Inches(0.65), Inches(10), Inches(0.5),
                 size=13, color=c, bold=True)

def slide_lineage_catalog():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Lineage y catálogo de datos", "M2")
    add_text(s, "Lineage = de dónde viene y a dónde va cada campo",
             Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             size=18, bold=True, color=NAVY)
    add_text(s, "Por qué importa",
             Inches(0.5), Inches(1.9), Inches(6), Inches(0.4),
             size=16, bold=True, color=CORAL)
    bullets(s, [
        "Impact analysis: si cambio esta columna, ¿qué rompe?",
        "Auditoría: ¿qué datos alimentaron este modelo en producción?",
        "RGPD/AI Act: trazabilidad obligatoria para decisiones automatizadas.",
        "Confianza: el consumidor sabe el origen del KPI.",
    ], Inches(0.7), Inches(2.4), Inches(6), Inches(3), size=14)

    add_text(s, "Herramientas",
             Inches(7), Inches(1.9), Inches(6), Inches(0.4),
             size=16, bold=True, color=CORAL)
    bullets(s, [
        ("OpenLineage", "estándar abierto, integra con Airflow, dbt"),
        ("DataHub",     "catálogo open source, LinkedIn"),
        ("Amundsen",    "search-first, Lyft"),
        ("Atlas",       "Apache, Hadoop ecosystem"),
        ("Unity Catalog", "Databricks, comercial"),
    ], Inches(7.2), Inches(2.4), Inches(6), Inches(3), size=14)

    add_round(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7), fill=CORAL)
    add_text(s, "Anti-pattern: 'la documentación está en Confluence' ≠ lineage",
             Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_dataops_antipatterns():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "DataOps anti-patterns que cuestan caro", "M2")
    items = [
        ("'Lo arreglamos en producción'", "El error de calidad se propaga al modelo y al dashboard"),
        ("Una sola tabla 'staging_final_v3'", "Sin versionado real, todo el mundo escribe encima"),
        ("Pipelines manuales", "Si el pipeline depende de que María ejecute el script, no es pipeline"),
        ("Validar al final", "El bug aparece en las predicciones, no en los datos"),
        ("Schema implícito", "Cambia un nombre de columna y todo cae en cascada"),
        ("Sin timezone", "Eventos de medianoche se duplican o desaparecen"),
        ("Sin contrato con upstream", "El proveedor cambia el formato sin avisar = silenciosa muerte"),
    ]
    bullets(s, items, Inches(0.5), Inches(1.3), Inches(12.5), Inches(5.5),
            size=15, bullet="✗")

def slide_lab1_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Lab 1 — DataOps reproducible (50 min)", "M2")
    add_round(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7), fill=CORAL)
    add_text(s, "Objetivo: pipeline de datos versionado, validado y reproducible",
             Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Pasos",
             Inches(0.5), Inches(2.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Inicializar repo con DVC apuntando a MinIO.",
        "Añadir dataset (UCI Adult Income).",
        "Stage 'validate' con Great Expectations (4 checks).",
        "Stage 'preprocess': encoding + split.",
        "Romper a propósito el dataset → ver el fallo.",
        "dvc repro y observar la ejecución parcial.",
    ], Inches(0.7), Inches(2.7), Inches(6), Inches(4), size=14)

    add_text(s, "Comprobaciones",
             Inches(7), Inches(2.2), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Existe data/raw/adult.csv.dvc en Git.",
        "El blob real está en MinIO bucket 'datasets'.",
        "ge_validate.py imprime ✓ y deja un Data Doc.",
        "dvc.lock tiene los hashes de cada stage.",
        "git commit del cambio NO incluye el csv.",
    ], Inches(7.2), Inches(2.7), Inches(6), Inches(4), size=14)

def slide_lab1_demo():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Demo en vivo: dvc.yaml + ge", "M2")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.6), fill=NAVY)
    code = ("# 1. inicializar\n"
            "$ dvc init\n"
            "$ dvc remote add -d minio s3://datasets\n"
            "$ dvc remote modify minio endpointurl http://localhost:9000\n\n"
            "# 2. añadir dataset\n"
            "$ dvc add data/raw/adult.csv     # crea adult.csv.dvc\n"
            "$ git add data/raw/adult.csv.dvc data/raw/.gitignore\n\n"
            "# 3. dvc.yaml\n"
            "stages:\n"
            "  validate:\n"
            "    cmd: python src/ge_validate.py data/raw/adult.csv\n"
            "    deps: [data/raw/adult.csv, src/ge_validate.py]\n"
            "  preprocess:\n"
            "    cmd: python src/preprocess.py\n"
            "    deps: [data/raw/adult.csv, src/preprocess.py]\n"
            "    outs: [data/processed/train.parquet, data/processed/test.parquet]\n\n"
            "# 4. ejecutar\n"
            "$ dvc repro\n"
            "$ dvc push  # sube blobs a MinIO")
    add_text(s, code, Inches(0.7), Inches(1.4), Inches(12), Inches(5.4),
             size=13, color=CYAN, font="Consolas")

def slide_m2_recap():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Recap Módulo 2", "M2")
    bullets(s, [
        "DataOps no es 'usar Airflow', es disciplina + automatización + calidad.",
        "Versiona datos con DVC (o LakeFS/Delta) para reproducir y auditar.",
        "Calidad medible en 7 dimensiones; tests con GE/Pandera/Soda.",
        "Lineage no es opcional cuando hay decisiones automatizadas.",
        "El pipeline acaba en la API/feature store, no en un notebook.",
    ], Inches(0.5), Inches(1.5), Inches(12), Inches(5), size=20)

# ============================================================================
# MÓDULO 3 — MLflow
# ============================================================================
def slide_ml_lifecycle():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Ciclo de vida de un modelo ML", "M3")
    stages = [
        ("Hipótesis",   CYAN),
        ("EDA",         CORAL),
        ("Preparar",    GREEN),
        ("Entrenar",    NAVY),
        ("Evaluar",     RGBColor(0xA0,0x55,0xF7)),
        ("Registrar",   RGBColor(0xF5,0x9E,0x0B)),
        ("Desplegar",   RGBColor(0xEC,0x49,0x99)),
        ("Monitor",     RGBColor(0x06,0xB6,0xD4)),
    ]
    for i, (n, c) in enumerate(stages):
        col = i % 4; row = i // 4
        x = Inches(0.5 + col*3.2); y = Inches(1.5 + row*1.5)
        add_round(s, x, y, Inches(3), Inches(1.2), fill=c)
        add_text(s, n, x, y, Inches(3), Inches(1.2),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # flecha circular implícita
    add_round(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), fill=BG, line=NAVY)
    add_text(s, "Es un ciclo: monitor → nueva hipótesis. Si no se cierra, el modelo decae.",
             Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5),
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_what_is_tracking():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "¿Qué es 'experiment tracking'?", "M3")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.9), fill=CYAN)
    add_text(s, "Registrar los metadatos de cada ejecución de entrenamiento de forma comparable",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.9),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Qué se registra",
             Inches(0.5), Inches(2.4), Inches(6), Inches(0.4),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        ("Parámetros", "lr, depth, semilla, versión data"),
        ("Métricas",   "f1, auc, recall por clase, MAE"),
        ("Artefactos", "modelo, gráficos, matriz confusión"),
        ("Tags",       "owner, dataset_version, hardware"),
        ("Entorno",    "requirements.txt, python version"),
        ("Source",     "git commit hash"),
    ], Inches(0.7), Inches(2.9), Inches(6), Inches(4), size=15)

    add_text(s, "Para qué sirve",
             Inches(7), Inches(2.4), Inches(6), Inches(0.4),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Comparar 200 runs sin abrir un Excel.",
        "Reproducir cualquier resultado del último año.",
        "Auditar: qué entrenó este modelo en producción.",
        "Compartir progreso con el equipo sin notebooks.",
        "Promocionar el mejor modelo a Staging/Prod.",
    ], Inches(7.2), Inches(2.9), Inches(6), Inches(4), size=15)

def slide_mlflow_components():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "MLflow: 4 componentes", "M3")
    comps = [
        ("Tracking", "API + UI para loggear params, métricas y artefactos", CYAN),
        ("Projects", "Empaqueta código en formato reproducible (MLproject yaml)", CORAL),
        ("Models",   "Formato unificado para guardar modelos (flavors)", GREEN),
        ("Registry", "Hub central de modelos con versiones y stages", NAVY),
    ]
    for i, (n, d, c) in enumerate(comps):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(1.4 + row*2.7)
        add_round(s, x, y, Inches(6.2), Inches(2.4), fill=WHITE, line=c)
        add_rect(s, x, y, Inches(6.2), Inches(0.7), fill=c)
        add_text(s, n, x, y, Inches(6.2), Inches(0.7),
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, d, x+Inches(0.3), y+Inches(0.9), Inches(5.6), Inches(1.5),
                 size=16, color=INK)

def slide_mlflow_tracking_code():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "MLflow Tracking: API mínima", "M3")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), fill=NAVY)
    code = ("import mlflow, mlflow.sklearn\n"
            "from sklearn.ensemble import RandomForestClassifier\n"
            "from sklearn.metrics import f1_score\n\n"
            "mlflow.set_tracking_uri('http://localhost:5000')\n"
            "mlflow.set_experiment('income-classifier')\n\n"
            "with mlflow.start_run(run_name='rf-baseline') as run:\n"
            "    mlflow.log_params({'n_estimators': 200, 'max_depth': 8})\n"
            "    mlflow.set_tag('dataset_version', 'v1.2')\n"
            "    mlflow.set_tag('owner', 'jpicon')\n\n"
            "    model = RandomForestClassifier(n_estimators=200, max_depth=8)\n"
            "    model.fit(X_train, y_train)\n\n"
            "    f1 = f1_score(y_test, model.predict(X_test))\n"
            "    mlflow.log_metric('f1', f1)\n"
            "    mlflow.sklearn.log_model(\n"
            "        model, name='model',\n"
            "        signature=mlflow.models.infer_signature(X_train, y_train),\n"
            "        input_example=X_train.head(3))")
    add_text(s, code, Inches(0.7), Inches(1.4), Inches(12), Inches(5.3),
             size=13, color=CYAN, font="Consolas")

def slide_mlflow_registry():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "MLflow Model Registry: el flujo", "M3")
    states = [("None", GREY), ("Staging", CORAL), ("Production", GREEN), ("Archived", NAVY)]
    x = Inches(0.5)
    for i, (n, c) in enumerate(states):
        add_round(s, x, Inches(1.5), Inches(2.5), Inches(1.0), fill=c)
        add_text(s, n, x, Inches(1.5), Inches(2.5), Inches(1.0),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(states) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     x+Inches(2.55), Inches(1.8), Inches(0.6), Inches(0.4))
            arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background()
        x += Inches(3.15)

    add_text(s, "Flujo recomendado",
             Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Cada run produce una versión candidata (v1, v2, v3…).",
        "Tras evaluación automática, la mejor pasa a Staging.",
        "Tras revisión humana + tests de integración → Production.",
        "Cuando se reemplaza, la anterior pasa a Archived (no se borra).",
        "Toda transición queda auditada (quién, cuándo, comentario).",
    ], Inches(0.7), Inches(3.6), Inches(12), Inches(3), size=15)

def slide_alternatives_tracking():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Alternativas a MLflow", "M3")
    rows = [
        ("Weights & Biases", "SaaS, UI brillante, integraciones LLMs", "Equipos con presupuesto, foco en investigación", CORAL),
        ("Neptune.ai",       "Tracking + monitoring, on-prem opcional", "Equipos enterprise con compliance",         CYAN),
        ("Comet ML",         "Foco en colaboración y reportes",         "Equipos product-data",                       GREEN),
        ("Aim",              "Open source ligero, self-host",           "Equipos pequeños sin K8s",                   NAVY),
        ("ClearML",          "Plataforma full open source",             "Quieren todo en uno + remote agents",        RGBColor(0xA0,0x55,0xF7)),
    ]
    for i, (n, d, u, c) in enumerate(rows):
        y = Inches(1.3 + i*1.05)
        add_round(s, Inches(0.5), y, Inches(2.6), Inches(0.9), fill=c)
        add_text(s, n, Inches(0.5), y, Inches(2.6), Inches(0.9),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(3.2), y, Inches(9.6), Inches(0.9), fill=WHITE, line=c)
        add_text(s, d, Inches(3.4), y+Inches(0.05), Inches(9.2), Inches(0.4),
                 size=13, color=INK)
        add_text(s, f"Ideal: {u}", Inches(3.4), y+Inches(0.45), Inches(9.2), Inches(0.4),
                 size=12, color=c, bold=True)

def slide_reproducibility():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Reproducibilidad: los 3 pilares", "M3")
    pillars = [
        ("Código", "Git commit hash + tag\nNo se entrena 'desde main'", CYAN),
        ("Datos",  "Hash del dataset (DVC)\nVersión registrada en el run", CORAL),
        ("Entorno","Imagen Docker + lockfile\nRequirements pinneados", GREEN),
    ]
    for i, (n, d, c) in enumerate(pillars):
        x = Inches(0.5 + i*4.27)
        add_round(s, x, Inches(1.5), Inches(4.1), Inches(3.5), fill=c)
        add_text(s, n, x, Inches(1.5), Inches(4.1), Inches(0.8),
                 size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, d, x+Inches(0.2), Inches(2.5), Inches(3.7), Inches(2.3),
                 size=16, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), fill=NAVY)
    add_text(s, "Si falta uno de los tres, no es reproducible. Punto.",
             Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4),
             size=22, bold=True, color=CORAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_signatures():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Model Signature e input_example", "M3")
    add_text(s, "Por qué",
             Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "Documenta el contrato del modelo (qué entra, qué sale).",
        "Permite validar inputs en serving automáticamente.",
        "Detecta cambios de schema entre runs.",
        "Integración directa con FastAPI / serving MLflow.",
    ], Inches(0.7), Inches(1.7), Inches(6), Inches(3), size=15)

    add_text(s, "Cómo",
             Inches(7), Inches(1.2), Inches(6), Inches(0.5),
             size=20, bold=True, color=NAVY)
    add_round(s, Inches(7), Inches(1.7), Inches(6), Inches(5), fill=NAVY)
    code = ("from mlflow.models import infer_signature\n\n"
            "sig = infer_signature(X_train, y_pred)\n\n"
            "mlflow.sklearn.log_model(\n"
            "    sk_model=model,\n"
            "    name='model',\n"
            "    signature=sig,\n"
            "    input_example=X_train.head(3),\n"
            "    pip_requirements=[\n"
            "        'scikit-learn==1.5.2',\n"
            "        'pandas==2.2.2'\n"
            "    ]\n"
            ")")
    add_text(s, code, Inches(7.2), Inches(1.85), Inches(5.7), Inches(4.7),
             size=13, color=CYAN, font="Consolas")

def slide_lab2_intro():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Lab 2 — MLflow Tracking + Registry (50 min)", "M3")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7), fill=GREEN)
    add_text(s, "Objetivo: experimentos reproducibles y un modelo en Staging",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Pasos",
             Inches(0.5), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Conectar al MLflow Tracking del docker-compose.",
        "Entrenar 3 modelos: LogReg, RandomForest, XGBoost.",
        "Loggear params, metrics, signature, input_example.",
        "Comparar runs en la UI (gráfico paralelo).",
        "Registrar el mejor → versión 1.",
        "Promover a Staging desde la API.",
    ], Inches(0.7), Inches(2.6), Inches(6), Inches(4), size=14)

    add_text(s, "Comprobaciones",
             Inches(7), Inches(2.1), Inches(6), Inches(0.5),
             size=18, bold=True, color=NAVY)
    bullets(s, [
        "Experimento 'income-classifier' tiene ≥3 runs.",
        "Cada run tiene tag 'dataset_version' = output del Lab 1.",
        "El run ganador tiene un modelo con signature.",
        "El Registry muestra 'income-clf v1 → Staging'.",
        "Puedes recuperar el modelo con 'models:/income-clf/Staging'.",
    ], Inches(7.2), Inches(2.6), Inches(6), Inches(4), size=14)

def slide_tagging_practice():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Buenas prácticas de tagging", "M3")
    add_text(s, "Tags imprescindibles en CADA run",
             Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             size=20, bold=True, color=NAVY)
    rows = [
        ("git_commit",        "hash del commit que entrenó"),
        ("dataset_version",   "ej. dvc-hash o tag semántico v1.2"),
        ("data_window",       "rango temporal de los datos"),
        ("owner",             "humano responsable"),
        ("ticket",            "JIRA/Linear que motiva el experimento"),
        ("hardware",          "cpu/gpu, memoria — relevante para coste"),
        ("framework_version", "scikit, torch, etc."),
        ("base_model",        "si es fine-tuning"),
    ]
    for i, (k, v) in enumerate(rows):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col*6.4); y = Inches(2.0 + row*0.7)
        add_round(s, x, y, Inches(2.5), Inches(0.6), fill=NAVY)
        add_text(s, k, x, y, Inches(2.5), Inches(0.6),
                 size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, v, x+Inches(2.6), y, Inches(3.7), Inches(0.6),
                 size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)

def slide_tracking_antipatterns():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Anti-patterns en experimentación", "M3")
    items = [
        "Loggear solo métricas: pierdes el modelo si lo necesitas reentrenar.",
        "Un experimento gigante con 5000 runs sin filtros: imposible navegar.",
        "Run sin git_commit: no sabes qué código produjo el resultado.",
        "Run sin dataset_version: no sabes qué dato produjo el resultado.",
        "Renombrar runs a mano: el equipo pierde la trazabilidad.",
        "'Promocionar' un modelo copiando un .pkl: bypassea el Registry.",
        "Confiar en una sola métrica: F1 alto puede ocultar mal recall por clase.",
    ]
    bullets(s, items, Inches(0.5), Inches(1.3), Inches(12.5), Inches(5.5),
            size=16, bullet="✗")

def slide_d1_close():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    header(s, "Cierre Día 1 + reto", "fin")
    add_text(s, "Lo que tenemos al final del día",
             Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             size=20, bold=True, color=NAVY)
    bullets(s, [
        "Datos versionados en MinIO con DVC y validados con Great Expectations.",
        "3 experimentos comparables en MLflow.",
        "Un modelo en Staging listo para servirse mañana.",
    ], Inches(0.7), Inches(1.8), Inches(12), Inches(2), size=17)

    add_round(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.2), fill=CORAL)
    add_text(s, "Reto del día (opcional)",
             Inches(0.7), Inches(4.1), Inches(12), Inches(0.5),
             size=18, bold=True, color=WHITE)
    add_text(s, ("Añade una expectation propia (ej: 'native-country' debe estar en una lista cerrada).\n"
                 "Entrena un cuarto modelo con búsqueda de hiperparámetros y vence al actual.\n"
                 "Subimos resultados al canal del curso. Mañana revisamos en 5 min."),
             Inches(0.7), Inches(4.6), Inches(12), Inches(1.6),
             size=15, color=WHITE)

# Construir
slides_funcs = [
    slide_cover,                                     # 1
    slide_about_speaker,                             # 2
    slide_agenda_d1,                                 # 3
    slide_objectives,                                # 4
    lambda: slide_module_divider(1, "Introducción a MLOps & DataOps", "Por qué existe la disciplina y qué problema resuelve", CYAN),  # 5
    slide_why_mlops,                                 # 6
    slide_hidden_debt,                               # 7
    slide_three_dimensions,                          # 8
    slide_devops_vs,                                 # 9
    slide_crisp_dm,                                  # 10
    lambda: slide_maturity("L0", "Manual / Notebook-driven",
        ["Notebooks como entregable final.",
         "Despliegue manual, sin tests de modelo.",
         "Sin versionado de datos ni modelos.",
         "Comunicación humana entre data y SW eng."],
        ["No se reproduce nada que no sea trivial.",
         "Mantener el modelo ocupa más que crearlo.",
         "Auditoría imposible (¿quién entrenó esto?)."],
        GREY),                                       # 11
    lambda: slide_maturity("L1", "ML pipeline automatizado",
        ["Pipeline reproducible orquestado.",
         "Datos y modelo versionados.",
         "Reentrenamiento programado o por trigger.",
         "Hay Model Registry con stages."],
        ["Promover modelo todavía requiere gente.",
         "CI/CD del código, no del pipeline ML completo.",
         "El registro existe pero el rollback no se prueba."],
        CORAL),                                      # 12
    lambda: slide_maturity("L2", "CI/CD pipelines completo",
        ["Pipeline ML completo en CI/CD.",
         "Test, train, eval, promote y deploy automatizados.",
         "Continuous Training (CT) con triggers de drift.",
         "Observabilidad completa modelo + datos + servicio."],
        ["Coste real: requiere plataforma y equipo dedicado.",
         "No tiene sentido sin volumen y madurez de datos.",
         "Sobre-ingeniería en proyectos de un solo modelo."],
        GREEN),                                      # 13
    slide_roles,                                     # 14
    slide_modern_stack,                              # 15
    slide_diagnostic_exercise,                       # 16
    lambda: slide_module_divider(2, "DataOps", "Pipelines, versionado y calidad de datos", CORAL),  # 17
    slide_dataops_manifesto,                         # 18
    slide_data_pipeline,                             # 19
    slide_why_data_versioning,                       # 20
    slide_dvc_alternatives,                          # 21
    slide_dvc_internals,                             # 22
    slide_data_quality_dimensions,                   # 23
    slide_great_expectations,                        # 24
    slide_pandera_soda,                              # 25
    slide_orchestrators,                             # 26
    slide_lineage_catalog,                           # 27
    slide_dataops_antipatterns,                      # 28
    slide_lab1_intro,                                # 29
    slide_lab1_demo,                                 # 30
    slide_m2_recap,                                  # 31
    lambda: slide_module_divider(3, "Experiment Tracking", "MLflow, Model Registry y reproducibilidad", GREEN),  # 32
    slide_ml_lifecycle,                              # 33
    slide_what_is_tracking,                          # 34
    slide_mlflow_components,                         # 35
    slide_mlflow_tracking_code,                      # 36
    slide_mlflow_registry,                           # 37
    slide_alternatives_tracking,                     # 38
    slide_reproducibility,                           # 39
    slide_signatures,                                # 40
    slide_lab2_intro,                                # 41
    slide_tagging_practice,                          # 42
    slide_tracking_antipatterns,                     # 43
    slide_d1_close,                                  # 44
]

for f in slides_funcs:
    f()

# Pagenum a todas excepto la portada
total = len(prs.slides)
for i, slide in enumerate(prs.slides):
    if i == 0:
        continue
    footer_pagenum(slide, i+1, total)

out = "/home/jmpicon/Documentos/Anbam/slides/dia1_mlops_dataops.pptx"
prs.save(out)
print(f"OK · {total} slides · {out}")
